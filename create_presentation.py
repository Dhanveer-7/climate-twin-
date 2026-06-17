import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Color Scheme (ISRO & Climatology inspired theme)
    DARK_BLUE = RGBColor(11, 19, 43)       # Deep Space Blue
    BRIGHT_ORANGE = RGBColor(242, 100, 25) # ISRO Rocket Orange
    LIGHT_BLUE = RGBColor(0, 180, 216)     # Accent Cyan
    WHITE = RGBColor(255, 255, 255)
    DARK_GREY = RGBColor(40, 40, 40)
    LIGHT_GREY = RGBColor(240, 242, 245)   # Light Clean Page BG
    
    # Standard blank layout (index 6 is blank)
    blank_slide_layout = prs.slide_layouts[6]
    
    def apply_solid_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, dark_mode=False):
        # Header text frame
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BRIGHT_ORANGE if dark_mode else DARK_BLUE
        
        # Add decorative colored accent line under title
        line_color = LIGHT_BLUE if dark_mode else BRIGHT_ORANGE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = line_color
        shape.line.fill.background()

    def add_footer(slide, dark_mode=False):
        # Footer text box for team and event identification
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4))
        tf = txBox.text_frame
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = "ISRO Bharatiya Antariksh Hackathon 2026  |  Designed by Team DhruvaX  |  Digital Twin of India's Climate"
        p.font.name = 'Arial'
        p.font.size = Pt(10)
        p.font.color.rgb = LIGHT_BLUE if dark_mode else RGBColor(120, 120, 120)
        p.alignment = PP_ALIGN.RIGHT

    # -------------------------------------------------------------------------
    # Slide 1: Title Slide (Dark Space Theme)
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide1, DARK_BLUE)
    
    # Hackathon Banner text
    tx_banner = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(2.2))
    tf1 = tx_banner.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "BHARATIYA ANTARIKSH HACKATHON 2026"
    p1.font.name = 'Trebuchet MS'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = BRIGHT_ORANGE
    
    p2 = tf1.add_paragraph()
    p2.text = "AI-Powered Digital Twin of India's Climate"
    p2.font.name = 'Trebuchet MS'
    p2.font.size = Pt(44)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.space_before = Pt(10)
    
    p3 = tf1.add_paragraph()
    p3.text = "Spatial-Temporal Anomaly Analytics & Vulnerability Simulator"
    p3.font.name = 'Arial'
    p3.font.size = Pt(18)
    p3.font.color.rgb = LIGHT_BLUE
    p3.space_before = Pt(10)
    
    # Details Box
    tx_details = slide1.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.5))
    tfd = tx_details.text_frame
    tfd.word_wrap = True
    
    p_team = tfd.paragraphs[0]
    p_team.text = "Team Name: DhruvaX"
    p_team.font.name = 'Trebuchet MS'
    p_team.font.size = Pt(20)
    p_team.font.bold = True
    p_team.font.color.rgb = WHITE
    
    p_leader = tfd.add_paragraph()
    p_leader.text = "Team Leader: Dhanveer M"
    p_leader.font.name = 'Arial'
    p_leader.font.size = Pt(16)
    p_leader.font.color.rgb = LIGHT_BLUE
    p_leader.space_before = Pt(6)
    
    p_prob = tfd.add_paragraph()
    p_prob.text = "Problem Statement: AI-Powered Digital Twin of India's Climate using India's National Data"
    p_prob.font.name = 'Arial'
    p_prob.font.size = Pt(14)
    p_prob.font.color.rgb = RGBColor(190, 200, 220)
    p_prob.space_before = Pt(15)

    # -------------------------------------------------------------------------
    # Slide 2: Team Members (Light Layout)
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide2, LIGHT_GREY)
    add_header(slide2, "Team Members")
    add_footer(slide2)
    
    members = [
        {"role": "Team Leader", "name": "Dhanveer M", "college": "College / Institution Name"},
        {"role": "Team Member 1", "name": "Jovin Joshua", "college": "College / Institution Name"},
        {"role": "Team Member 2", "name": "Elamaran B", "college": "College / Institution Name"},
        {"role": "Team Member 3", "name": "Athish K.B", "college": "College / Institution Name"}
    ]
    
    positions = [
        (Inches(1.0), Inches(1.8)),
        (Inches(7.0), Inches(1.8)),
        (Inches(1.0), Inches(4.3)),
        (Inches(7.0), Inches(4.3))
    ]
    
    for idx, member in enumerate(members):
        pos_x, pos_y = positions[idx]
        
        # Outer Card Box
        shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pos_y, Inches(5.333), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_BLUE
        shape.line.width = Pt(1.5)
        
        tf_card = shape.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.25)
        tf_card.margin_top = Inches(0.2)
        
        p_role = tf_card.paragraphs[0]
        p_role.text = member["role"]
        p_role.font.name = 'Trebuchet MS'
        p_role.font.size = Pt(14)
        p_role.font.bold = True
        p_role.font.color.rgb = BRIGHT_ORANGE
        
        p_name = tf_card.add_paragraph()
        p_name.text = f"Name: {member['name']}"
        p_name.font.name = 'Trebuchet MS'
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = DARK_BLUE
        p_name.space_before = Pt(8)
        
        p_coll = tf_card.add_paragraph()
        p_coll.text = f"College: {member['college']}"
        p_coll.font.name = 'Arial'
        p_coll.font.size = Pt(12)
        p_coll.font.color.rgb = DARK_GREY
        p_coll.space_before = Pt(6)

    # -------------------------------------------------------------------------
    # Slide 3: Opportunity should be able to explain the following (Light Layout)
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide3, LIGHT_GREY)
    add_header(slide3, "Opportunity & Uniqueness")
    add_footer(slide3)
    
    cols = [
        {
            "title": "How different is it from existing ideas?",
            "bullets": [
                "Unified Meteorological Analysis: Integrates spatial temperature, rainfall, and multi-state anomalies into a single operating dashboard.",
                "High-Fidelity Virtual Twin: Computes realistic rolling seasonality instead of simple linear models."
            ]
        },
        {
            "title": "How will it solve the problem?",
            "bullets": [
                "Translates complex ISRO and national climatology data into daily interactive visual maps.",
                "Enables localized multi-horizon forecasting (7-day and 30-day horizons) to let administrators plan ahead."
            ]
        },
        {
            "title": "USP of the Proposed Solution",
            "bullets": [
                "Dynamic What-If Sandbox: Run warming simulations (+1°C to +3°C) and rainfall delta checks to see immediate sectoral effects.",
                "AI-Driven Warning Cards: Automatically alerts about heat hazards and crop droughts."
            ]
        }
    ]
    
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    start_x = Inches(0.8)
    
    for idx, col in enumerate(cols):
        pos_x = start_x + idx * (col_width + col_gap)
        
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, Inches(1.8), col_width, Inches(4.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BRIGHT_ORANGE if idx == 2 else LIGHT_BLUE
        shape.line.width = Pt(1.5)
        
        tf_card = shape.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.2)
        tf_card.margin_right = Inches(0.2)
        tf_card.margin_top = Inches(0.2)
        
        p_title = tf_card.paragraphs[0]
        p_title.text = col["title"]
        p_title.font.name = 'Trebuchet MS'
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK_BLUE
        
        for bullet in col["bullets"]:
            parts = bullet.split(":", 1)
            p_bullet = tf_card.add_paragraph()
            p_bullet.font.name = 'Trebuchet MS'
            p_bullet.space_before = Pt(12)
            
            if len(parts) > 1:
                p_bullet.text = "• " + parts[0].strip() + ":"
                p_bullet.font.bold = True
                p_bullet.font.size = Pt(12)
                p_bullet.font.color.rgb = BRIGHT_ORANGE if idx == 2 else DARK_BLUE
                
                p_desc = tf_card.add_paragraph()
                p_desc.text = parts[1].strip()
                p_desc.font.name = 'Arial'
                p_desc.font.size = Pt(11)
                p_desc.font.color.rgb = DARK_GREY
                p_desc.space_before = Pt(2)
                p_desc.margin_left = Inches(0.15)
            else:
                p_bullet.text = "• " + bullet.strip()
                p_bullet.font.size = Pt(11)
                p_bullet.font.color.rgb = DARK_GREY

    # -------------------------------------------------------------------------
    # Slide 4: List of Features Offered by the Solution (Light Layout)
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide4, LIGHT_GREY)
    add_header(slide4, "List of Features Offered by the Solution")
    add_footer(slide4)
    
    features = [
        {"title": "Meteorological Overview Dashboard", "desc": "Clean temporal analytics of state-wise temperature and monsoon rainfall anomaly indices."},
        {"title": "GIS Operations Room Map", "desc": "Interactive geospatial maps overlaying Z-scores and departure maps for state-by-state assessment."},
        {"title": "Recursive AI Forecasting Engine", "desc": "Pre-trained XGBoost and RF models for multi-step 7-day and 30-day climate forecasting."},
        {"title": "What-If Stress Sandbox", "desc": "Simulate and stress-test regional climate response to temperature anomalies and rainfall variances."},
        {"title": "Active Warning Cards & Alarms", "desc": "Automated alarm banners indicating severe heat hazards, agricultural drought, and flood levels."},
        {"title": "Offline Local DB Cache", "desc": "SQLite database caching 220,000+ daily meteorological records to run efficiently offline."}
    ]
    
    feat_width = Inches(5.6)
    feat_height = Inches(1.3)
    
    for idx, feat in enumerate(features):
        col_idx = idx % 2
        row_idx = idx // 2
        
        pos_x = Inches(0.8) if col_idx == 0 else Inches(6.8)
        pos_y = Inches(1.8) + row_idx * Inches(1.6)
        
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pos_y, feat_width, feat_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_BLUE
        shape.line.width = Pt(1)
        
        tf_feat = shape.text_frame
        tf_feat.word_wrap = True
        tf_feat.margin_left = Inches(0.2)
        tf_feat.margin_top = Inches(0.15)
        
        p_title = tf_feat.paragraphs[0]
        p_title.text = f"{idx+1}. {feat['title']}"
        p_title.font.name = 'Trebuchet MS'
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = BRIGHT_ORANGE
        
        p_desc = tf_feat.add_paragraph()
        p_desc.text = feat["desc"]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GREY
        p_desc.space_before = Pt(4)

    # -------------------------------------------------------------------------
    # Slide 5: Process Flow Diagram or Use-case Diagram (Light Layout)
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide5, LIGHT_GREY)
    add_header(slide5, "Process Flow & Data Operations")
    add_footer(slide5)
    
    steps = [
        {"title": "1. Data Ingestion", "desc": "SQLite database loading 220,000+ daily meteorology records (2010-2026)"},
        {"title": "2. Feature Engineering", "desc": "Calculates rolling seasonality & multi-step forecast indicators"},
        {"title": "3. Simulation Engine", "desc": "Applies user-selected offsets and recursive AI forecasting models"},
        {"title": "4. Actionable Alerts", "desc": "Renders Folium GIS Maps, threat scores, and warning card banners"}
    ]
    
    box_width = Inches(2.5)
    box_height = Inches(2.2)
    start_x = Inches(0.6)
    pos_y = Inches(2.6)
    
    for idx, step in enumerate(steps):
        pos_x = start_x + idx * Inches(3.2)
        
        shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pos_y, box_width, box_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = DARK_BLUE
        shape.line.width = Pt(2)
        
        tf_box = shape.text_frame
        tf_box.word_wrap = True
        tf_box.margin_left = Inches(0.15)
        tf_box.margin_top = Inches(0.2)
        
        p_title = tf_box.paragraphs[0]
        p_title.text = step["title"]
        p_title.font.name = 'Trebuchet MS'
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = BRIGHT_ORANGE
        p_title.alignment = PP_ALIGN.CENTER
        
        p_desc = tf_box.add_paragraph()
        p_desc.text = step["desc"]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GREY
        p_desc.space_before = Pt(10)
        p_desc.alignment = PP_ALIGN.CENTER
        
        # Draw Arrow connecting boxes
        if idx < 3:
            arrow_x = pos_x + box_width + Inches(0.1)
            arrow_y = pos_y + Inches(0.9)
            arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.5), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_BLUE
            arrow.line.fill.background()

    # -------------------------------------------------------------------------
    # Slide 6: Wireframes/Mock diagrams of the proposed solution (Light Layout)
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide6, LIGHT_GREY)
    add_header(slide6, "Wireframes & Live Application Interface")
    add_footer(slide6)
    
    # Left column - App specs
    txBox = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_link = tf.paragraphs[0]
    p_link.text = "Deployed Application URL:"
    p_link.font.name = 'Trebuchet MS'
    p_link.font.size = Pt(16)
    p_link.font.bold = True
    p_link.font.color.rgb = DARK_BLUE
    
    p_url = tf.add_paragraph()
    p_url.text = "https://dhanveer-7-climate-twin--app-eodlti.streamlit.app/"
    p_url.font.name = 'Courier New'
    p_url.font.size = Pt(13)
    p_url.font.bold = True
    p_url.font.color.rgb = BRIGHT_ORANGE
    p_url.space_before = Pt(5)
    
    p_feat_title = tf.add_paragraph()
    p_feat_title.text = "Key UI Sandboxes and Layouts:"
    p_feat_title.font.name = 'Trebuchet MS'
    p_feat_title.font.size = Pt(16)
    p_feat_title.font.bold = True
    p_feat_title.font.color.rgb = DARK_BLUE
    p_feat_title.space_before = Pt(20)
    
    uis = [
        "1. Operations Room Map: Fully interactive state-wise climate anomaly overlay.",
        "2. ML Predictions Board: Side-by-side comparison of XGBoost vs RF forecasts.",
        "3. Scenario Simulator: User-interactive gauge sliders to alter ΔT and ΔR.",
        "4. Warning Card System: Instantly generates threat levels for farming/health sectors."
    ]
    for ui in uis:
        p_ui = tf.add_paragraph()
        p_ui.text = "• " + ui.split(":")[0] + ":"
        p_ui.font.name = 'Trebuchet MS'
        p_ui.font.bold = True
        p_ui.font.size = Pt(12)
        p_ui.font.color.rgb = LIGHT_BLUE
        p_ui.space_before = Pt(8)
        
        p_ui_desc = tf.add_paragraph()
        p_ui_desc.text = ui.split(":")[1].strip()
        p_ui_desc.font.name = 'Arial'
        p_ui_desc.font.size = Pt(11)
        p_ui_desc.font.color.rgb = DARK_GREY
        p_ui_desc.space_before = Pt(1)
        p_ui_desc.margin_left = Inches(0.15)
        
    # Right column - Graphic visual card placeholder
    shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_BLUE
    shape.line.width = Pt(1.5)
    
    tf_right = shape.text_frame
    tf_right.word_wrap = True
    tf_right.margin_left = Inches(0.3)
    tf_right.margin_top = Inches(0.4)
    
    p_right_t = tf_right.paragraphs[0]
    p_right_t.text = "Interactive Screen Walkthrough"
    p_right_t.font.name = 'Trebuchet MS'
    p_right_t.font.size = Pt(18)
    p_right_t.font.bold = True
    p_right_t.font.color.rgb = BRIGHT_ORANGE
    
    p_right_d = tf_right.add_paragraph()
    p_right_d.text = (
        "The digital twin has been deployed to Streamlit Community Cloud and is accessible "
        "globally. It features a space-inspired Dark Blue Operations Dashboard styling tailored for the ISRO hackathon theme.\n\n"
        "Features responsive sidebars with Click-Only Dropdowns to improve the user experience on tablets and laptops, "
        "completely avoiding default keyboard blinking cursor boxes."
    )
    p_right_d.font.name = 'Arial'
    p_right_d.font.size = Pt(13)
    p_right_d.font.color.rgb = DARK_BLUE
    p_right_d.space_before = Pt(15)

    # -------------------------------------------------------------------------
    # Slide 7: Architecture diagram of the proposed solution (Light Layout)
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide7, LIGHT_GREY)
    add_header(slide7, "System Architecture")
    add_footer(slide7)
    
    layers = [
        {"title": "Presentation Layer (Dashboard UI)", "desc": "Streamlit Frontend + Custom CSS Theme + Folium GIS Maps + Plotly Gauge Charts"},
        {"title": "Simulation Engine (What-If Sandbox)", "desc": "What-If Simulator + Dynamic Stress Calculators + Sectoral Impact Warning Cards"},
        {"title": "Machine Learning Layer (AI Forecasting)", "desc": "Pre-trained RF Regressors & XGBoost Regressors (Multi-step recursive 7-day/30-day forecast loops)"},
        {"title": "Data Storage Layer (Relational DB)", "desc": "Local SQLite Database containing 223,524 daily historical state records (2010-2026)"}
    ]
    
    for idx, layer in enumerate(layers):
        pos_y = Inches(1.8) + idx * Inches(1.2)
        shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), pos_y, Inches(11.333), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BRIGHT_ORANGE if idx == 0 else LIGHT_BLUE
        shape.line.width = Pt(1.5)
        
        tf_l = shape.text_frame
        tf_l.word_wrap = True
        tf_l.margin_left = Inches(0.2)
        tf_l.margin_top = Inches(0.15)
        
        p_title = tf_l.paragraphs[0]
        p_title.text = layer["title"]
        p_title.font.name = 'Trebuchet MS'
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = DARK_BLUE
        
        p_desc = tf_l.add_paragraph()
        p_desc.text = layer["desc"]
        p_desc.font.name = 'Arial'
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = DARK_GREY
        p_desc.space_before = Pt(3)

    # -------------------------------------------------------------------------
    # Slide 8: Technologies to be used in the solution (Light Layout)
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide8, LIGHT_GREY)
    add_header(slide8, "Technologies & Libraries")
    add_footer(slide8)
    
    techs = [
        {"cat": "Web & Dashboard", "items": ["Streamlit (v1.40+)", "HTML5 / Custom CSS (Operations Theme)", "Jinja2 Templates"]},
        {"cat": "Machine Learning", "items": ["XGBoost Regressor", "Scikit-Learn (Random Forest)", "Pickle Serialization"]},
        {"cat": "Data Storage", "items": ["SQLite3 Database Engine", "Pandas Climatology Pipelines", "NumPy Numerical Operations"]},
        {"cat": "Geospatial & Charts", "items": ["Folium Map Overlay", "Plotly Express", "Branca (HTML Map Rendering)"]}
    ]
    
    card_width = Inches(5.4)
    card_height = Inches(2.1)
    
    card_positions = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.2)),
        (Inches(6.8), Inches(4.2))
    ]
    
    for idx, tech in enumerate(techs):
        pos_x, pos_y = card_positions[idx]
        shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pos_y, card_width, card_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_BLUE
        shape.line.width = Pt(1.5)
        
        tf_t = shape.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = Inches(0.25)
        tf_t.margin_top = Inches(0.2)
        
        p_cat = tf_t.paragraphs[0]
        p_cat.text = tech["cat"]
        p_cat.font.name = 'Trebuchet MS'
        p_cat.font.size = Pt(16)
        p_cat.font.bold = True
        p_cat.font.color.rgb = BRIGHT_ORANGE
        
        for item in tech["items"]:
            p_item = tf_t.add_paragraph()
            p_item.text = f"• {item}"
            p_item.font.name = 'Arial'
            p_item.font.size = Pt(12)
            p_item.font.color.rgb = DARK_BLUE
            p_item.space_before = Pt(4)

    # -------------------------------------------------------------------------
    # Slide 9: Estimated implementation cost (Light Layout)
    # -------------------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide9, LIGHT_GREY)
    add_header(slide9, "Implementation Cost & Scalability Strategy")
    add_footer(slide9)
    
    tx_left = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_l = tx_left.text_frame
    tf_l.word_wrap = True
    
    p_l_title = tf_l.paragraphs[0]
    p_l_title.text = "Implementation Cost Assessment"
    p_l_title.font.name = 'Trebuchet MS'
    p_l_title.font.size = Pt(18)
    p_l_title.font.bold = True
    p_l_title.font.color.rgb = DARK_BLUE
    
    costs = [
        "Infrastructure Cost: Deployed on free cloud containers (Streamlit Cloud). Enterprise scaling costs < $10/month.",
        "Computation: Pre-trained regressors run prediction loops on CPU in milliseconds, bypassing GPU needs.",
        "Database Size: Relational database takes < 50MB, saving high database hosting budgets."
    ]
    for cost in costs:
        p_cost = tf_l.add_paragraph()
        p_cost.text = "• " + cost.split(":")[0] + ":"
        p_cost.font.name = 'Trebuchet MS'
        p_cost.font.bold = True
        p_cost.font.size = Pt(12)
        p_cost.font.color.rgb = BRIGHT_ORANGE
        p_cost.space_before = Pt(12)
        
        p_cost_desc = tf_l.add_paragraph()
        p_cost_desc.text = cost.split(":")[1].strip()
        p_cost_desc.font.name = 'Arial'
        p_cost_desc.font.size = Pt(11)
        p_cost_desc.font.color.rgb = DARK_GREY
        p_cost_desc.space_before = Pt(2)
        p_cost_desc.margin_left = Inches(0.15)
        
    tx_right = slide9.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_r = tx_right.text_frame
    tf_r.word_wrap = True
    
    p_r_title = tf_r.paragraphs[0]
    p_r_title.text = "Future Scalability Roadmap"
    p_r_title.font.name = 'Trebuchet MS'
    p_r_title.font.size = Pt(18)
    p_r_title.font.bold = True
    p_r_title.font.color.rgb = DARK_BLUE
    
    scales = [
        "INSAT Stream Sync: Ingesting satellite telemetry streams (monsoon clouds) dynamically.",
        "Ground Sensors (REST API): Connect state weather grids directly to relational SQL db.",
        "Deep Learning Models: Replacing regressors with ConvLSTMs to capture complex multi-layer meteorological movements."
    ]
    for scale in scales:
        p_scale = tf_r.add_paragraph()
        p_scale.text = "• " + scale.split(":")[0] + ":"
        p_scale.font.name = 'Trebuchet MS'
        p_scale.font.bold = True
        p_scale.font.size = Pt(12)
        p_scale.font.color.rgb = LIGHT_BLUE
        p_scale.space_before = Pt(12)
        
        p_scale_desc = tf_r.add_paragraph()
        p_scale_desc.text = scale.split(":")[1].strip()
        p_scale_desc.font.name = 'Arial'
        p_scale_desc.font.size = Pt(11)
        p_scale_desc.font.color.rgb = DARK_GREY
        p_scale_desc.space_before = Pt(2)
        p_scale_desc.margin_left = Inches(0.15)

    # -------------------------------------------------------------------------
    # Slide 10: Thank You (Dark Space Theme)
    # -------------------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    apply_solid_background(slide10, DARK_BLUE)
    
    tx_thank = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf10 = tx_thank.text_frame
    tf10.word_wrap = True
    
    p_thank = tf10.paragraphs[0]
    p_thank.text = "THANK YOU"
    p_thank.font.name = 'Trebuchet MS'
    p_thank.font.size = Pt(64)
    p_thank.font.bold = True
    p_thank.font.color.rgb = BRIGHT_ORANGE
    p_thank.alignment = PP_ALIGN.CENTER
    
    p_sub = tf10.add_paragraph()
    p_sub.text = "Bharatiya Antariksh Hackathon 2026"
    p_sub.font.name = 'Trebuchet MS'
    p_sub.font.size = Pt(24)
    p_sub.font.bold = True
    p_sub.font.color.rgb = WHITE
    p_sub.space_before = Pt(20)
    p_sub.alignment = PP_ALIGN.CENTER
    
    p_brand = tf10.add_paragraph()
    p_brand.text = "Presented by Team DhruvaX  |  Empowering Climate Resilience with Digital Twins"
    p_brand.font.name = 'Trebuchet MS'
    p_brand.font.size = Pt(16)
    p_brand.font.color.rgb = LIGHT_BLUE
    p_brand.space_before = Pt(15)
    p_brand.alignment = PP_ALIGN.CENTER

    output_path = "dhruvax_climate_twin_presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
